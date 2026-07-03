"""Local RAG engine with document loading, index validation, and vector search."""
import hashlib
import json
import pickle
import re
from datetime import datetime, timezone
from pathlib import Path

import numpy as np
import requests

from config import EMBED_MODEL, LMSTUDIO_BASE_URL, VECTOR_STORE_DIR


CHUNK_SIZE = 500
CHUNK_OVERLAP = 50
INDEX_SCHEMA_VERSION = 2
SUPPORTED_EXTENSIONS = {".txt", ".md", ".pdf", ".pptx"}


class RagError(RuntimeError):
    """Base exception for the RAG layer."""


class EmbeddingServiceError(RagError):
    """Embedding service is unavailable or returned an invalid response."""


class IndexStaleError(RagError):
    """Index no longer matches the current documents or configuration."""


def _infer_title(text: str, fallback: str) -> str:
    """Infer a short display title from a text block."""
    for line in text.splitlines():
        line = re.sub(r"^#{1,6}\s*", "", line).strip(" -|\t")
        if 2 <= len(line) <= 80:
            return line
    return fallback


def _table_to_markdown(table: list[list]) -> str:
    """Convert PDF/PPT tables into searchable Markdown tables."""
    rows = []
    for row in table or []:
        cells = [str(cell or "").replace("\n", " ").strip() for cell in row]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    header = rows[0]
    body = rows[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in body)
    return "\n".join(lines)


def _split_slide_text(text: str, source: str) -> list[dict]:
    """Split preprocessed slide text into sections when slide markers exist."""
    pattern = re.compile(
        r"^(?:---\s*幻灯片\s+(\d+)\s*---|##\s*幻灯片\s+(\d+)(?:[：:].*)?)\s*$",
        re.MULTILINE,
    )
    matches = list(pattern.finditer(text))
    if not matches:
        return [{"text": text, "source": source, "title": _infer_title(text, Path(source).stem)}]

    sections = []
    for i, match in enumerate(matches):
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        section_text = text[start:end].strip()
        if section_text:
            slide = int(match.group(1) or match.group(2))
            sections.append({
                "text": section_text,
                "source": source,
                "title": _infer_title(section_text, f"Slide {slide}"),
                "slide": slide,
            })
    return sections


def load_text_file(filepath: str) -> list[dict]:
    path = Path(filepath)
    text = path.read_text(encoding="utf-8")
    return _split_slide_text(text, path.name)


def load_pdf_file(filepath: str) -> list[dict]:
    """Load PDF pages, preserving tables when pdfplumber is installed."""
    path = Path(filepath)
    sections = []
    try:
        import pdfplumber

        with pdfplumber.open(path) as pdf:
            for page_number, page in enumerate(pdf.pages, 1):
                text = (page.extract_text() or "").strip()
                tables = []
                for table in page.extract_tables() or []:
                    markdown = _table_to_markdown(table)
                    if markdown:
                        tables.append(markdown)
                combined = text
                if tables:
                    combined += "\n\n### Tables\n\n" + "\n\n".join(tables)
                if combined.strip():
                    sections.append({
                        "text": combined.strip(),
                        "source": path.name,
                        "title": _infer_title(text, f"Page {page_number}"),
                        "page": page_number,
                    })
        return sections
    except ImportError:
        pass
    except Exception as exc:
        print(f"[WARN] pdfplumber failed, falling back to PyPDF2: {path.name}: {exc}")

    try:
        from PyPDF2 import PdfReader
    except ImportError as exc:
        raise RagError("Missing PDF parser dependency. Install pdfplumber or PyPDF2.") from exc

    reader = PdfReader(path)
    for page_number, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            sections.append({
                "text": text,
                "source": path.name,
                "title": _infer_title(text, f"Page {page_number}"),
                "page": page_number,
            })
    return sections


def load_pptx_file(filepath: str) -> list[dict]:
    """Load PPTX slides while preserving slide titles and Markdown tables."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RagError("Missing PPT parser dependency. Install python-pptx.") from exc

    path = Path(filepath)
    presentation = Presentation(path)
    sections = []
    for slide_number, slide in enumerate(presentation.slides, 1):
        text_blocks = []
        tables = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                block = "\n".join(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
                if block:
                    text_blocks.append(block)
            if getattr(shape, "has_table", False):
                table = [[cell.text for cell in row.cells] for row in shape.table.rows]
                markdown = _table_to_markdown(table)
                if markdown:
                    tables.append(markdown)
        combined = "\n\n".join(text_blocks)
        if tables:
            combined += "\n\n### Tables\n\n" + "\n\n".join(tables)
        if combined.strip():
            sections.append({
                "text": combined.strip(),
                "source": path.name,
                "title": _infer_title(combined, f"Slide {slide_number}"),
                "slide": slide_number,
            })
    return sections


def load_folder(folder: str | Path) -> list[dict]:
    docs = []
    folder_path = Path(folder)
    if not folder_path.exists():
        return docs
    for filepath in sorted(folder_path.iterdir()):
        if filepath.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        try:
            if filepath.suffix.lower() == ".pdf":
                loaded = load_pdf_file(str(filepath))
            elif filepath.suffix.lower() == ".pptx":
                loaded = load_pptx_file(str(filepath))
            else:
                loaded = load_text_file(str(filepath))
            for document in loaded:
                document["source_folder"] = folder_path.name
                document["source_path"] = str(filepath)
            docs.extend(loaded)
        except Exception as exc:
            print(f"[WARN] Failed to load {filepath.name}: {exc}")
    return docs


def load_folders(folders: str | Path | list[str | Path] | tuple[str | Path, ...]) -> list[dict]:
    if isinstance(folders, (str, Path)):
        return load_folder(folders)

    docs = []
    for folder in folders:
        docs.extend(load_folder(folder))
    return docs


def _split_long_text(text: str, limit: int) -> list[str]:
    """Split long paragraphs on sentence-like boundaries."""
    if len(text) <= limit:
        return [text]
    sentences = re.split(r"(?<=[。！？；.!?;])\s*|\n+", text)
    parts, current = [], ""
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if len(current) + len(sentence) + 1 <= limit:
            current = (current + "\n" + sentence).strip()
        else:
            if current:
                parts.append(current)
            if len(sentence) > limit:
                parts.extend(sentence[i:i + limit] for i in range(0, len(sentence), limit))
                current = ""
            else:
                current = sentence
    if current:
        parts.append(current)
    return parts


def chunk_document(document: dict, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[dict]:
    """Chunk a document section while keeping page, slide, and title metadata."""
    paragraphs = []
    for paragraph in re.split(r"\n\s*\n", document["text"]):
        paragraph = paragraph.strip()
        if paragraph:
            paragraphs.extend(_split_long_text(paragraph, chunk_size))

    chunks, current = [], ""
    for paragraph in paragraphs:
        if len(current) + len(paragraph) + 2 <= chunk_size:
            current = (current + "\n\n" + paragraph).strip()
            continue
        if current:
            chunks.append(current)
            overlap_text = current[-overlap:] if overlap else ""
            current = (overlap_text + "\n\n" + paragraph).strip()
        else:
            current = paragraph
    if current:
        chunks.append(current)

    result = []
    for index, text in enumerate(chunks):
        item = {key: value for key, value in document.items() if key != "text"}
        item.update({
            "id": (
                f"{document.get('source_folder', '')}/{document['source']}"
                f"#{document.get('page', document.get('slide', 0))}#{index}"
            ),
            "text": text,
            "chunk_index": index,
        })
        result.append(item)
    return result


def _embed_batch(texts: list[str]) -> list[list[float]]:
    """Call the embedding service and raise clear errors on failure."""
    try:
        response = requests.post(
            f"{LMSTUDIO_BASE_URL}/v1/embeddings",
            json={"model": EMBED_MODEL, "input": texts},
            timeout=120,
        )
    except requests.ConnectionError as exc:
        raise EmbeddingServiceError(
            f"Cannot connect to the local embedding service ({LMSTUDIO_BASE_URL}). "
            f"Start LM Studio and load {EMBED_MODEL}."
        ) from exc
    except requests.Timeout as exc:
        raise EmbeddingServiceError("Local embedding service request timed out.") from exc
    except requests.RequestException as exc:
        raise EmbeddingServiceError(f"Embedding service request failed: {exc}") from exc

    if not response.ok:
        detail = response.text[:300].strip()
        raise EmbeddingServiceError(f"Embedding service returned HTTP {response.status_code}: {detail}")
    try:
        data = response.json()["data"]
        return [item["embedding"] for item in data]
    except (KeyError, TypeError, ValueError) as exc:
        raise EmbeddingServiceError("Embedding service returned an unexpected payload.") from exc


def embed_chunks(chunks: list[dict], on_progress=None) -> list[list[float]]:
    all_embeddings = []
    texts = [chunk["text"] for chunk in chunks]
    for i in range(0, len(texts), 10):
        all_embeddings.extend(_embed_batch(texts[i:i + 10]))
        if on_progress:
            on_progress(min(i + 10, len(texts)), len(texts))
    return all_embeddings


def _source_manifest(folder: str | Path | list[str | Path] | tuple[str | Path, ...]) -> list[dict]:
    if not isinstance(folder, (str, Path)):
        manifest = []
        for item in folder:
            manifest.extend(_source_manifest(item))
        return sorted(manifest, key=lambda entry: (entry["folder"], entry["name"]))

    folder_path = Path(folder)
    if not folder_path.exists():
        return []
    manifest = []
    for filepath in sorted(folder_path.iterdir()):
        if filepath.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        digest = hashlib.sha256(filepath.read_bytes()).hexdigest()
        stat = filepath.stat()
        manifest.append({
            "folder": folder_path.name,
            "name": filepath.name,
            "size": stat.st_size,
            "sha256": digest,
        })
    return manifest


def _manifest_fingerprint(manifest: list[dict]) -> str:
    payload = json.dumps(manifest, ensure_ascii=False, sort_keys=True).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def format_citation(chunk: dict) -> str:
    source = chunk.get("source", "Unknown source")
    if chunk.get("source_folder"):
        source = f"{chunk['source_folder']}/{source}"
    parts = [source]
    if chunk.get("page"):
        parts.append(f"Page {chunk['page']}")
    if chunk.get("slide"):
        parts.append(f"Slide {chunk['slide']}")
    if chunk.get("title"):
        parts.append(chunk["title"])
    return " · ".join(parts)


def _serialize_docs_folder(docs_folder: str | Path | list[str | Path] | tuple[str | Path, ...]) -> str | list[str]:
    if isinstance(docs_folder, (str, Path)):
        return str(Path(docs_folder).resolve())
    return [str(Path(folder).resolve()) for folder in docs_folder]


class VectorStore:
    def __init__(self, store_dir: str | Path = None):
        self.store_dir = Path(store_dir or VECTOR_STORE_DIR)
        self.store_dir.mkdir(parents=True, exist_ok=True)
        self.chunks: list[dict] = []
        self.embeddings: np.ndarray | None = None
        self.metadata: dict = {}
        self._loaded = False
        self._stale_reason = ""

    def build(self, docs_folder: str | Path | list[str | Path] | tuple[str | Path, ...]) -> bool:
        docs = load_folders(docs_folder)
        if not docs:
            raise RagError(f"No supported documents found in: {docs_folder}")

        chunks = []
        for document in docs:
            chunks.extend(chunk_document(document))
        if not chunks:
            raise RagError("Documents were read, but no indexable text was produced.")

        embeddings = embed_chunks(chunks)
        manifest = _source_manifest(docs_folder)
        self.chunks = chunks
        self.embeddings = np.asarray(embeddings, dtype=np.float32)
        self.metadata = {
            "schema_version": INDEX_SCHEMA_VERSION,
            "built_at": datetime.now(timezone.utc).isoformat(),
            "embedding_model": EMBED_MODEL,
            "embedding_dimension": int(self.embeddings.shape[1]),
            "chunk_size": CHUNK_SIZE,
            "chunk_overlap": CHUNK_OVERLAP,
            "docs_folder": _serialize_docs_folder(docs_folder),
            "source_manifest": manifest,
            "source_fingerprint": _manifest_fingerprint(manifest),
        }
        self._loaded = True
        self._stale_reason = ""
        self._save()
        return True

    def _save(self):
        with (self.store_dir / "index.pkl").open("wb") as file:
            pickle.dump({"chunks": self.chunks, "embeddings": self.embeddings, "metadata": self.metadata}, file)
        (self.store_dir / "index_metadata.json").write_text(
            json.dumps(self.metadata, ensure_ascii=False, indent=2), encoding="utf-8"
        )

    def load(self, docs_folder: str | Path | list[str | Path] | tuple[str | Path, ...] = None) -> bool:
        index_path = self.store_dir / "index.pkl"
        if not index_path.exists():
            return False
        with index_path.open("rb") as file:
            data = pickle.load(file)
        self.chunks = data.get("chunks", [])
        self.embeddings = data.get("embeddings")
        self.metadata = data.get("metadata", {})
        self._loaded = True
        self._validate(docs_folder)
        return True

    def _validate(self, docs_folder: str | Path | list[str | Path] | tuple[str | Path, ...] = None):
        reasons = []
        if self.metadata.get("schema_version") != INDEX_SCHEMA_VERSION:
            reasons.append("index schema changed")
        if self.metadata.get("embedding_model") != EMBED_MODEL:
            reasons.append("embedding model changed")
        if self.metadata.get("chunk_size") != CHUNK_SIZE or self.metadata.get("chunk_overlap") != CHUNK_OVERLAP:
            reasons.append("chunking configuration changed")
        if docs_folder is not None:
            folders = [docs_folder] if isinstance(docs_folder, (str, Path)) else list(docs_folder)
            missing = [str(folder) for folder in folders if not Path(folder).exists()]
            if missing:
                reasons.append(f"document folder missing: {', '.join(missing)}")
            else:
                manifest = _source_manifest(docs_folder)
                fingerprint = _manifest_fingerprint(manifest)
                if self.metadata.get("source_fingerprint") != fingerprint:
                    reasons.append("source files were added, removed, or modified")

        self._stale_reason = "; ".join(reasons)

    @property
    def is_loaded(self):
        return self._loaded and self.embeddings is not None and bool(self.chunks)

    @property
    def is_stale(self):
        return bool(self._stale_reason)

    @property
    def stale_reason(self):
        return self._stale_reason

    @property
    def chunk_count(self):
        return len(self.chunks)

    @property
    def embedding_dimension(self):
        if self.embeddings is None or self.embeddings.ndim != 2:
            return None
        return int(self.embeddings.shape[1])

    def search(self, query: str, top_k: int = 5) -> list[dict]:
        if not self.is_loaded:
            raise RagError("Knowledge index has not been built yet.")
        if self.is_stale:
            raise IndexStaleError(f"Index needs to be rebuilt: {self.stale_reason}")
        if not query.strip():
            return []

        query_vec = np.asarray(_embed_batch([query])[0], dtype=np.float32)
        query_length = np.linalg.norm(query_vec)
        embedding_lengths = np.linalg.norm(self.embeddings, axis=1, keepdims=True)
        if query_length == 0 or np.any(embedding_lengths == 0):
            raise RagError("Embedding returned a zero vector, so similarity cannot be calculated.")
        scores = np.dot(self.embeddings / embedding_lengths, query_vec / query_length)

        boost_terms = [
            "setup", "configuration", "pricing", "safety", "workflow", "support",
            "integration", "deployment", "training", "security", "privacy", "compliance",
        ]
        query_lower = query.lower()
        for index, chunk in enumerate(self.chunks):
            text_lower = chunk["text"].lower()
            bonus = sum(0.04 for term in boost_terms if term in query_lower and term in text_lower)
            scores[index] += min(bonus, 0.12)

        top_indices = np.argsort(scores)[::-1][:max(top_k * 4, 20)]
        results, seen_sources = [], {}
        for index in top_indices:
            score = float(scores[index])
            if score < 0.2:
                continue
            chunk = self.chunks[index]
            source_key = (chunk.get("source"), chunk.get("page"), chunk.get("slide"))
            count = seen_sources.get(source_key, 0)
            if count >= 2:
                continue
            seen_sources[source_key] = count + 1
            results.append({**chunk, "score": round(score, 4), "citation": format_citation(chunk)})
            if len(results) >= top_k:
                break
        return results


_store: VectorStore | None = None


def get_store(docs_folder: str | Path | list[str | Path] | tuple[str | Path, ...] = None) -> VectorStore:
    global _store
    if _store is None:
        _store = VectorStore()
        _store.load(docs_folder)
    elif docs_folder is not None:
        _store._validate(docs_folder)
    return _store


def rebuild_store(docs_folder: str | Path | list[str | Path] | tuple[str | Path, ...]) -> VectorStore:
    global _store
    new_store = VectorStore()
    new_store.build(docs_folder)
    _store = new_store
    return _store
