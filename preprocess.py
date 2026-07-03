"""
Document preprocessing script. Converts PPT/PPTX into structured Markdown
for the local knowledge base.

Usage:
  python preprocess.py <document-folder-path>

Examples:
  python preprocess.py "./sample_docs"
  python preprocess.py "./sample_docs/demo_deck.pptx"

Output:
  Creates a converted/ folder next to the input and stores extracted Markdown.
"""
import sys
import os
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Missing python-pptx. Installing...")
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation


def _table_to_markdown(shape) -> str:
    rows = []
    for row in shape.table.rows:
        cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n".join(lines)


def extract_text_from_pptx(filepath: str) -> str:
    """Extract slide titles, text, and Markdown tables from PPTX."""
    prs = Presentation(filepath)
    all_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_title = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            slide_title = slide.shapes.title.text.strip()
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_text.append(text)
            if shape.has_table:
                table_text = _table_to_markdown(shape)
                if table_text:
                    slide_text.append("### Tables\n\n" + table_text)
        if slide_text:
            heading = f"## Slide {slide_num}"
            if slide_title:
                heading += f": {slide_title}"
            all_text.append(heading + "\n\n" + "\n\n".join(slide_text))
    return "\n\n".join(all_text)


def process_file(filepath: Path, output_dir: Path) -> str:
    """Process one file and return the output path."""
    ext = filepath.suffix.lower()
    output_name = filepath.stem + ".md"

    if ext in (".pptx", ".ppt"):
        try:
            text = extract_text_from_pptx(str(filepath))
        except Exception as e:
            print(f"  [WARN] PPT extraction failed: {e}")
            return None

        if not text.strip():
            print("  [WARN] No text was extracted. The deck may be image-only.")
            return None

        output_path = output_dir / output_name
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        return str(output_path)

    elif ext == ".pdf":
        import shutil
        dest = output_dir / filepath.name
        shutil.copy2(str(filepath), str(dest))
        return str(dest)

    elif ext in (".docx", ".txt", ".md", ".csv"):
        import shutil
        dest = output_dir / filepath.name
        shutil.copy2(str(filepath), str(dest))
        return str(dest)

    else:
        print(f"  [SKIP] Unsupported format: {ext}")
        return None


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        return

    input_path = Path(sys.argv[1])

    if not input_path.exists():
        print(f"[ERROR] Path does not exist: {input_path}")
        return

    if input_path.is_file():
        output_dir = input_path.parent / "converted"
    else:
        output_dir = input_path / "converted"

    output_dir.mkdir(exist_ok=True)

    supported = {".pptx", ".ppt", ".pdf", ".docx", ".txt", ".md", ".csv"}
    files = []
    if input_path.is_file():
        files = [input_path]
    else:
        for root, _, filenames in os.walk(input_path):
            for f in filenames:
                fp = Path(root) / f
                if fp.suffix.lower() in supported:
                    files.append(fp)

    if not files:
        print(f"[ERROR] No supported files found ({', '.join(supported)})")
        return

    print(f"[INFO] Found {len(files)} files. Processing...\n")

    success = 0
    for fp in files:
        print(f"  Processing: {fp.name}")
        result = process_file(fp, output_dir)
        if result:
            print(f"    [OK] -> {Path(result).name}")
            success += 1

    print(f"\n{'='*50}")
    print(f"[OK] Done. Processed {success}/{len(files)} files.")
    print(f"[INFO] Output directory: {output_dir}")
    print(f"\nNext: point LOCAL_RAG_DOCS_FOLDER at '{output_dir}' or copy files into sample_docs/.")


if __name__ == "__main__":
    main()
